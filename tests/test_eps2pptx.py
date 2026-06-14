import io
import tempfile
import unittest
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

FIXTURES = Path(__file__).parent / "fixtures"


class ConvertTests(unittest.TestCase):
    def _convert(self, name):
        import eps2pptx
        with tempfile.TemporaryDirectory() as d:
            dst = Path(d) / "out.pptx"
            status = eps2pptx.convert_eps_to_pptx(FIXTURES / name, dst)
            return dst.read_bytes(), status

    def test_grid_fixture_yields_nine_shapes(self):
        raw, status = self._convert("grid_3x3.eps")
        zf = zipfile.ZipFile(io.BytesIO(raw))
        slide = zf.read("ppt/slides/slide1.xml").decode("utf-8")
        self.assertEqual(slide.count("<a:custGeom>"), 9)
        self.assertIn("9", status)

    def test_all_parts_well_formed(self):
        raw, _ = self._convert("grid_3x3.eps")
        zf = zipfile.ZipFile(io.BytesIO(raw))
        for name in zf.namelist():
            if name.endswith(".xml") or name.endswith(".rels"):
                ET.fromstring(zf.read(name))

    @unittest.skipUnless(__import__("importlib").util.find_spec("pptx"),
                         "python-pptx not installed")
    def test_opens_in_python_pptx_with_nine_shapes(self):
        from pptx import Presentation
        with tempfile.TemporaryDirectory() as d:
            dst = Path(d) / "out.pptx"
            import eps2pptx
            eps2pptx.convert_eps_to_pptx(FIXTURES / "grid_3x3.eps", dst)
            prs = Presentation(str(dst))
            shapes = list(prs.slides[0].shapes)
            self.assertEqual(len(shapes), 9)


class ConvertIconsToPptxTests(unittest.TestCase):
    def _icons(self, n):
        # Each icon: one tiny square path fragment with its own bbox.
        icons = []
        for i in range(n):
            frag = f'<path d="M0 0 L10 0 L10 10 L0 10 Z" fill="#{i:02x}0000"/>'
            icons.append(([frag], (0.0, 0.0, 10.0, 10.0)))
        return icons

    def test_one_slide_per_icon(self):
        import eps2pptx
        with tempfile.TemporaryDirectory() as d:
            dst = Path(d) / "deck.pptx"
            status = eps2pptx.convert_icons_to_pptx(self._icons(4), dst)
            self.assertIn("4", status)
            zf = zipfile.ZipFile(io.BytesIO(dst.read_bytes()))
            slides = [n for n in zf.namelist()
                      if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
            self.assertEqual(len(slides), 4)
            for name in zf.namelist():
                if name.endswith(".xml") or name.endswith(".rels"):
                    ET.fromstring(zf.read(name))

    @unittest.skipUnless(__import__("importlib").util.find_spec("pptx"),
                         "python-pptx not installed")
    def test_deck_opens_in_python_pptx(self):
        from pptx import Presentation
        import eps2pptx
        with tempfile.TemporaryDirectory() as d:
            dst = Path(d) / "deck.pptx"
            eps2pptx.convert_icons_to_pptx(self._icons(3), dst)
            prs = Presentation(str(dst))
            self.assertEqual(len(prs.slides), 3)
